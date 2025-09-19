import json
import os
from pptx import Presentation
from pptx.util import Inches
import sys

def main():
    """
    Main function that processes a stringified JSON input to create PowerPoint slides.

    Expected JSON format:
    {
        "slides": [
            {
                "variant": "Cover.pptx",
                "content": {
                    "title": "My Presentation",
                    "subtitle": "A great presentation",
                    "date": "2024-01-01"
                }
            },
            ...
        ]
    }
    """
    try:
        # Read JSON input from stdin
        json_input = sys.stdin.read().strip()

        if not json_input:
            print("Error: No JSON input provided", file=sys.stderr)
            sys.exit(1)

        # Parse JSON
        data = json.loads(json_input)

        if 'slides' not in data:
            print("Error: JSON must contain 'slides' array", file=sys.stderr)
            sys.exit(1)

        # Get the directory containing slide templates
        script_dir = os.path.dirname(os.path.abspath(__file__))
        slides_dir = os.path.join(script_dir, '..', 'slides')

        print(f"Script directory: {script_dir}", file=sys.stderr)
        print(f"Slides directory: {slides_dir}", file=sys.stderr)
        print(f"Slides directory exists: {os.path.exists(slides_dir)}", file=sys.stderr)

        # Create a new presentation
        prs = Presentation()

        # Process each slide
        print(f"Processing {len(data['slides'])} slides", file=sys.stderr)
        for i, slide_data in enumerate(data['slides']):
            variant_name = slide_data.get('variant')
            content = slide_data.get('content', {})

            print(f"Slide {i+1}: variant={variant_name}, content keys={list(content.keys())}", file=sys.stderr)

            if not variant_name:
                print(f"Warning: Skipping slide without variant name", file=sys.stderr)
                continue

            # Find the variant file
            variant_path = os.path.join(slides_dir, variant_name)
            print(f"Looking for variant file: {variant_path}", file=sys.stderr)

            if not os.path.exists(variant_path):
                print(f"Warning: Variant file not found: {variant_path}", file=sys.stderr)
                continue

            # Load the template presentation
            try:
                template_prs = Presentation(variant_path)

                # Copy the first slide from template (assuming it's the slide we want)
                if len(template_prs.slides) > 0:
                    template_slide = template_prs.slides[0]

                    # Add a new slide to our presentation
                    slide_layout = prs.slide_layouts[0]  # Use blank layout
                    new_slide = prs.slides.add_slide(slide_layout)

                    # Copy slide content
                    copy_slide_content(template_slide, new_slide, content)

            except Exception as e:
                print(f"Error processing variant {variant_name}: {str(e)}", file=sys.stderr)
                continue

        # delete any existing output file
        output_path = os.path.join(script_dir, '..', 'output', 'GeneratedPresentation.pptx')
        if os.path.exists(output_path):
            os.remove(output_path)

        # Save the presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)

        print(f"Presentation saved to: {output_path}")

    except json.JSONDecodeError as e:
        print(f"Error parsing JSON: {str(e)}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)

def copy_slide_content(template_slide, new_slide, content):
    """
    Copy content from template slide to new slide and replace text based on content mapping.
    """
    # Copy all shapes from template to new slide
    for shape in template_slide.shapes:
        if hasattr(shape, 'text'):
            # This is a text shape
            original_text = shape.text

            # Try to replace text based on content mapping
            replaced_text = replace_text_content(original_text, content)

            # Add text box to new slide
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = replaced_text

            # Copy text formatting if possible
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                source_para = shape.text_frame.paragraphs[0]
                target_para = text_frame.paragraphs[0]

                if source_para.font.size:
                    target_para.font.size = source_para.font.size
                if source_para.font.bold is not None:
                    target_para.font.bold = source_para.font.bold
                if source_para.font.italic is not None:
                    target_para.font.italic = source_para.font.italic

        elif hasattr(shape, 'shape_type'):
            # This is another type of shape (image, etc.)
            # For now, we'll skip non-text shapes
            # In a more complete implementation, you'd copy these too
            pass

def replace_text_content(text, content):
    """
    Replace placeholder text with actual content based on the variants.json mapping.
    """
    # Load variants.json to understand the mapping
    script_dir = os.path.dirname(os.path.abspath(__file__))
    variants_path = os.path.join(script_dir, '..', 'variants.json')

    try:
        with open(variants_path, 'r') as f:
            variants = json.load(f)
    except FileNotFoundError:
        print(f"Warning: variants.json not found at {variants_path}", file=sys.stderr)
        return text

    # Find the variant that matches our content structure
    for variant in variants:
        variant_props = variant.get('properties', {})

        # Check if this variant matches our content keys
        if all(key in content for key in variant_props.keys()):
            # Replace text based on the variant mapping
            replaced_text = text

            for key, value in content.items():
                # Look for the placeholder value in the text
                if key in variant_props:
                    placeholder = variant_props[key]
                    if placeholder in replaced_text:
                        replaced_text = replaced_text.replace(placeholder, str(value))

            return replaced_text

    # If no specific variant found, try generic replacement
    # This is a fallback for cases where the variant mapping doesn't match exactly
    replaced_text = text
    for key, value in content.items():
        # Try to find and replace common patterns
        patterns = [
            f"{{{key}}}",
            f"{{{{{key}}}}}",
            f"<{key}>",
            f"<{key.upper()}>",
            key.upper(),
            key.title()
        ]

        for pattern in patterns:
            if pattern in replaced_text:
                replaced_text = replaced_text.replace(pattern, str(value))
                break

    return replaced_text

if __name__ == "__main__":
    main()
